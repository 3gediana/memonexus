"""知识库文本提取模块

对齐 Node.js 版 extractor.js 的 extractText 逻辑：
- PDF: PyMuPDF 先提文本层 → 文本少的页面走 RapidOCR-Paddle GPU 混合
- Office: python-docx / python-pptx / openpyxl
- 老格式(doc/xls/ppt): LibreOffice 转 PDF 再提取
- 图片: RapidOCR-Paddle GPU
- txt/md: 直接读取

OCR 使用 rapidocr-paddle（PaddlePaddle CUDA GPU 加速）
"""

import os
import re
import json
import subprocess
import logging
import numpy as np
from pathlib import Path

logger = logging.getLogger(__name__)

KB_DIR = str(Path(__file__).resolve().parent.parent.parent.parent / "知识库")

# Paddle cannot load models from paths with non-ASCII chars on Windows,
# so we symlink/junction them to a temp ASCII path at runtime.
_OCR_MODEL_LINK_DIR: str | None = None


def _ensure_ascii_model_dir() -> str:
    """Create a junction/symlink to the paddle model dir at an ASCII-only path.
    Paddle inference engine fails to open model files from Chinese-char paths on Windows.
    """
    global _OCR_MODEL_LINK_DIR
    if _OCR_MODEL_LINK_DIR is not None:
        return _OCR_MODEL_LINK_DIR

    import shutil, tempfile, platform

    src = os.path.join(KB_DIR, ".models", "paddleocr-models", "whl")

    # check if path is already ASCII-safe
    try:
        src.encode("ascii")
        _OCR_MODEL_LINK_DIR = src
        return src
    except UnicodeEncodeError:
        pass

    dst = os.path.join(tempfile.gettempdir(), "memonexus_paddle_models_whl")

    # try junction (Windows) or symlink (Unix)
    if platform.system() == "Windows":
        if os.path.exists(dst):
            # verify junction points to correct location
            try:
                # junction target check
                if os.path.realpath(dst) == os.path.realpath(src):
                    _OCR_MODEL_LINK_DIR = dst
                    return dst
            except Exception:
                pass
            try:
                os.rmdir(dst)
            except Exception:
                shutil.rmtree(dst, ignore_errors=True)

        try:
            import subprocess

            subprocess.run(
                ["cmd", "/c", "mklink", "/J", dst, src], capture_output=True, check=True
            )
            logger.info(f"Created junction {dst} -> {src}")
            _OCR_MODEL_LINK_DIR = dst
            return dst
        except Exception:
            pass
    else:
        if os.path.islink(dst) and os.path.realpath(dst) == os.path.realpath(src):
            _OCR_MODEL_LINK_DIR = dst
            return dst

    # fallback: copy files
    try:
        if os.path.exists(dst):
            shutil.rmtree(dst)
        shutil.copytree(src, dst)
        logger.info(f"Copied paddle models to {dst}")
        _OCR_MODEL_LINK_DIR = dst
        return dst
    except Exception as e:
        logger.error(f"Failed to prepare ASCII model dir: {e}")
        _OCR_MODEL_LINK_DIR = src
        return src


SUPPORTED_EXT = {
    ".pdf",
    ".docx",
    ".doc",
    ".xlsx",
    ".xls",
    ".pptx",
    ".ppt",
    ".txt",
    ".md",
    ".jpg",
    ".jpeg",
    ".png",
    ".bmp",
}

# RapidOCR-Paddle 全局单例
_ocr_instance = None


_ocr_init_failed = False


def _get_ocr():
    """获取 RapidOCR-Paddle 实例（单例懒加载，优先 mobile 模型+GPU，失败降级）"""
    global _ocr_instance, _ocr_init_failed
    if _ocr_instance is not None:
        return _ocr_instance
    if _ocr_init_failed:
        return None

    try:
        from rapidocr_paddle import RapidOCR

        model_dir = _ensure_ascii_model_dir()
        # prefer mobile models (small & fast), fallback to server models
        det_path = os.path.join(
            model_dir, "det", "ch", "ch_PP-OCRv4_det_infer_mobile_bak"
        )
        rec_path = os.path.join(
            model_dir, "rec", "ch", "ch_PP-OCRv4_rec_infer_mobile_bak"
        )
        cls_path = os.path.join(model_dir, "cls", "ch_ppocr_mobile_v2.0_cls_infer")

        # fallback to root-level models if mobile_bak missing
        if not os.path.isfile(os.path.join(det_path, "inference.pdmodel")):
            det_path = os.path.join(model_dir, "det")
        if not os.path.isfile(os.path.join(rec_path, "inference.pdmodel")):
            rec_path = os.path.join(model_dir, "rec")
        if not os.path.isfile(os.path.join(cls_path, "inference.pdmodel")):
            cls_path = os.path.join(model_dir, "cls")

        # try GPU first, fall back to CPU
        for use_gpu in [True, False]:
            try:
                _ocr_instance = RapidOCR(
                    det_use_gpu=use_gpu,
                    cls_use_gpu=use_gpu,
                    rec_use_gpu=use_gpu,
                    det_model_path=det_path,
                    rec_model_path=rec_path,
                    cls_model_path=cls_path,
                )
                # warm up with tiny image to catch allocation errors early
                import numpy as np

                _ocr_instance(np.zeros((32, 32, 3), dtype=np.uint8))
                logger.info(f"RapidOCR-Paddle initialized (gpu={use_gpu})")
                return _ocr_instance
            except Exception as e:
                logger.warning(f"RapidOCR init with gpu={use_gpu} failed: {e}")
                _ocr_instance = None
                continue

        if _ocr_instance is None:
            logger.error("RapidOCR-Paddle: all init attempts failed, OCR disabled")
            _ocr_init_failed = True
    except ImportError:
        logger.warning("rapidocr-paddle not installed, OCR unavailable")
        _ocr_init_failed = True
    except Exception as e:
        logger.warning(f"RapidOCR-Paddle unexpected error: {e}")
        _ocr_init_failed = True
    return None


def _ocr_image(img_array) -> str:
    """对单张图片做 OCR，返回识别文本"""
    ocr = _get_ocr()
    if ocr is None:
        return ""
    try:
        result, _ = ocr(img_array)
        if result:
            return "\n".join(item[1] for item in result)
    except Exception as e:
        logger.warning(f"OCR 识别失败: {e}")
    return ""


def _extract_pdf_text(file_path: str) -> str:
    """PDF 混合提取：先文本层，文本少的页面走 OCR（仅扫描版/图片页）"""
    import fitz

    try:
        doc = fitz.open(file_path)
        pdf_texts = []
        pages_need_ocr = []
        total_chars = 0

        for i in range(doc.page_count):
            page = doc[i]
            text = page.get_text("text")
            total_chars += len(text.strip())
            if len(text.strip()) > 100:
                pdf_texts.append(text)
            else:
                pages_need_ocr.append(i)
                pdf_texts.append("")

        # skip OCR if most pages have rich text (not a scanned PDF)
        avg_chars_per_page = total_chars / max(doc.page_count, 1)
        if pages_need_ocr and avg_chars_per_page > 500:
            logger.info(
                f"PDF {os.path.basename(file_path)}: avg {avg_chars_per_page:.0f} chars/page, "
                f"treating empty pages as blanks (not scanned)"
            )
            pages_need_ocr.clear()

        if pages_need_ocr:
            logger.info(
                f"PDF {os.path.basename(file_path)}: {len(pages_need_ocr)}/{doc.page_count} pages need OCR"
            )
            ocr = _get_ocr()
            if ocr is None:
                logger.warning("OCR unavailable, skipping OCR pages")
            else:
                for page_idx in pages_need_ocr:
                    page = doc[page_idx]
                    mat = fitz.Matrix(2.0, 2.0)
                    pix = page.get_pixmap(matrix=mat)
                    img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(
                        pix.height, pix.width, pix.n
                    )
                    if pix.n == 4:
                        img = img[:, :, :3]
                    try:
                        result, _ = ocr(img)
                        if result:
                            pdf_texts[page_idx] = "\n".join(item[1] for item in result)
                    except Exception as e:
                        logger.warning(f"page {page_idx + 1} OCR failed: {e}")

        doc.close()
        return "\n".join(t for t in pdf_texts if t.strip())

    except Exception as e:
        logger.error(f"PDF parse failed {file_path}: {e}")
        return ""


def _extract_docx_text(file_path: str) -> str:
    """提取 .docx 文本"""
    from docx import Document

    try:
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        logger.error(f"docx 解析失败 {file_path}: {e}")
        return ""


def _extract_pptx_text(file_path: str) -> str:
    """提取 .pptx 文本"""
    from pptx import Presentation

    try:
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text)
        return "\n".join(texts)
    except Exception as e:
        logger.error(f"pptx 解析失败 {file_path}: {e}")
        return ""


def _extract_xlsx_text(file_path: str) -> str:
    """提取 .xlsx 文本"""
    from openpyxl import load_workbook

    try:
        wb = load_workbook(file_path, read_only=True, data_only=True)
        texts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                row_text = " ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    texts.append(row_text)
        wb.close()
        return "\n".join(texts)
    except Exception as e:
        logger.error(f"xlsx 解析失败 {file_path}: {e}")
        return ""


def _convert_to_pdf(file_path: str) -> str:
    """通过 LibreOffice 将老格式文件转为 PDF"""
    libreoffice_paths = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "soffice",
    ]

    libreoffice_cmd = None
    for cmd in libreoffice_paths:
        if cmd == "soffice" or os.path.exists(cmd):
            libreoffice_cmd = cmd
            break

    if not libreoffice_cmd:
        logger.warning("LibreOffice 未安装，无法转换老格式文件")
        return ""

    tmp_dir = os.path.join(KB_DIR, ".scripts", ".tmp", "convert")
    os.makedirs(tmp_dir, exist_ok=True)

    try:
        result = subprocess.run(
            [
                libreoffice_cmd,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                tmp_dir,
                file_path,
            ],
            capture_output=True,
            text=True,
            timeout=120,
        )
        if result.returncode == 0:
            pdf_name = Path(file_path).stem + ".pdf"
            pdf_path = os.path.join(tmp_dir, pdf_name)
            if os.path.exists(pdf_path):
                return pdf_path
    except Exception as e:
        logger.warning(f"LibreOffice 转换失败 {file_path}: {e}")
    return ""


def extract_text(file_path: str) -> str:
    """统一文本提取入口

    Args:
        file_path: 文件绝对路径

    Returns:
        提取的纯文本
    """
    ext = Path(file_path).suffix.lower()

    if ext in (".txt", ".md"):
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
        except Exception as e:
            logger.error(f"读取文本文件失败 {file_path}: {e}")
            return ""

    if ext in (".jpg", ".jpeg", ".png", ".bmp"):
        return _ocr_image(file_path)

    if ext == ".pdf":
        return _extract_pdf_text(file_path)

    if ext == ".docx":
        return _extract_docx_text(file_path)

    if ext == ".pptx":
        return _extract_pptx_text(file_path)

    if ext in (".xlsx", ".xls"):
        if ext == ".xlsx":
            return _extract_xlsx_text(file_path)
        else:
            pdf_path = _convert_to_pdf(file_path)
            if pdf_path:
                text = _extract_pdf_text(pdf_path)
                try:
                    os.remove(pdf_path)
                except OSError:
                    pass
                return text
            return ""

    if ext in (".doc", ".ppt"):
        pdf_path = _convert_to_pdf(file_path)
        if pdf_path:
            text = _extract_pdf_text(pdf_path)
            try:
                os.remove(pdf_path)
            except OSError:
                pass
            return text
        return ""

    logger.warning(f"不支持的文件类型: {ext}")
    return ""
